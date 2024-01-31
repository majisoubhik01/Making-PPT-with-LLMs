import openai
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx import Presentation
from dotenv import load_dotenv
import os

load_dotenv()

openai.api_key = os.getenv('OPENAI_API_KEY')


TITLE_FONT_SIZE = Pt(32)
CONTENT_FONT_SIZE = Pt(16)

def create_slide_titles(topic, num_slides):
    prompt = f"Create {num_slides} short slide titles for the following topic: {topic}"

    completion = openai.ChatCompletion.create(
        model='gpt-3.5-turbo-1106',
        messages=[
            {'role': 'system', 'content': prompt}
        ],
        temperature=0.0,
        top_p=0.1,
        max_tokens=200,
        request_timeout=30
    )

    return completion.choices[0].message.content


def create_slide_content(slide_title):
    prompt = f"Generate content for the slide: {slide_title}. The content must be in medium worded paragraphs. Only return 2 paragraphs."

    completion = openai.ChatCompletion.create(
        model='gpt-3.5-turbo-1106',
        messages=[
            {'role': 'system', 'content': prompt}
        ],
        temperature=0.0,
        top_p=0.1,
        max_tokens=200,
        request_timeout=20
    )

    return completion.choices[0].message.content

def create_presentation(topic, slide_titles, slide_contents):
    powerpoint = Presentation()

    # 1st slide
    title_slide_layout = powerpoint.slide_layouts[0]
    content_slide_layout = powerpoint.slide_layouts[1]

    background_color = RGBColor(173,216,230)

    title_slide = powerpoint.slides.add_slide(title_slide_layout)

    title = title_slide.shapes.title
    title.text = topic
    title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title.text_frame.paragraphs[0].font.bold = True
    content = title_slide.placeholders[1]
    content.text = 'Created by LLM'

    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color

    # Rest of the slides
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        
        slide = powerpoint.slides.add_slide(content_slide_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        title.text_frame.paragraphs[0].font.bold = True

        content = slide.placeholders[1]
        content.text = slide_content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = CONTENT_FONT_SIZE

    powerpoint.save(f"powerpoint-{topic}.pptx")


def main():
    topic = "Genesis" #input("Enter the topic: ")

    num_slides = 1 #int(input("Enter the number of slides: "))

    slide_titles = create_slide_titles(topic, num_slides)
    print("Generated slide titles.")
    print(slide_titles)

    slide_contents = [create_slide_content(title) for title in slide_titles]
    print("Generated slide contents.")

    create_presentation(topic, slide_titles, slide_contents)
    print("Done!")

if __name__ == "__main__":
    main()

